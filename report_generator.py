import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import io
import textwrap
import re
import openai

openai_api_key = st.secrets["OPENAI_API_KEY"]

st.title("PPT 자동 생성기 (Streamlit)")

st.markdown('''
**❗️아래 예시처럼 반드시 번호+제목+내용 형태와 유사한 구조로 입력해야 합니다!**

예시: 사업계획인 경우
1. 사업의 필요성
사업의 필요성에 대한 내용

2. 사업의 개요
사업의 개요에 대한 내용

3. 기대효과
기대효과에 대한 내용
''')

st.markdown("문장을 복사해서 넣으면 요약해 PPT로 만들어줌 (만들고 싶은 페이지 수 선택)")

content = st.text_area("PPT로 만들고 싶은 내용을 입력하세요...", height=200)

col1, col2 = st.columns([1, 1])
with col1:
    page_count = st.number_input("페이지 수", min_value=1, value=3, step=1, key="page_count")
with col2:
    make_summary = st.button("요약 생성")

if make_summary:
    # 번호+제목+내용 구조로 슬라이드 분할
    items = re.findall(r'(\d+)\.\s*([^\n]+)\n([^\n]+(?:\n(?!\d+\.).+)*)', content, re.MULTILINE)
    slides_content = []
    client = openai.OpenAI(api_key=openai_api_key)
    for idx, title, item_content in items:
        prompt = f"다음 내용을 2~3줄의 bullet point로 요약해줘:\n{item_content.strip()}"
        try:
            response = client.chat.completions.create(
                model="gpt-3.5-turbo",
                messages=[{"role": "user", "content": prompt}],
                max_tokens=512,
                temperature=0.7,
            )
            summary = response.choices[0].message.content.strip()
        except Exception as e:
            summary = f"(요약 실패: {e})\n{item_content.strip()}"
        slides_content.append({'title': title.strip(), 'content': summary})

    # 슬라이드가 1개도 없으면 경고
    if not slides_content:
        st.error("❗️입력한 내용에서 번호(1. 2. 3. ...)로 구분된 슬라이드가 없습니다.\n\n아래 예시처럼 반드시 번호+제목+내용 구조로 입력해 주세요!\n\n예시:\n1. 사업의 필요성\n사업의 필요성에 대한 내용\n\n2. 사업의 개요\n사업의 개요에 대한 내용\n\n3. 기대효과\n기대효과에 대한 내용")
        st.session_state['grouped_slides'] = []
    else:
        def group_sections_by_page(sections, page_count):
            n = len(sections)
            base = n // page_count
            extra = n % page_count
            grouped = []
            idx = 0
            for i in range(page_count):
                count = base + (1 if i < extra else 0)
                grouped.append(sections[idx:idx+count])
                idx += count
            return grouped

        grouped_slides = group_sections_by_page(slides_content, int(page_count))
        st.session_state['grouped_slides'] = grouped_slides

if 'grouped_slides' in st.session_state and st.session_state['grouped_slides']:
    st.markdown("**페이지별 내용 미리보기 및 수정**")
    edited_grouped_slides = []
    for page_idx, group in enumerate(st.session_state['grouped_slides']):
        st.markdown(f"---\n### 페이지 {page_idx+1}")
        edited_group = []
        for slide_idx, slide in enumerate(group):
            title = st.text_input(f"페이지 {page_idx+1} - 슬라이드 {slide_idx+1} 제목", value=slide['title'], key=f"title_{page_idx}_{slide_idx}")
            body = st.text_area(f"페이지 {page_idx+1} - 슬라이드 {slide_idx+1} 내용", value=slide['content'], key=f"body_{page_idx}_{slide_idx}")
            edited_group.append({'title': title, 'content': body})
        edited_grouped_slides.append(edited_group)

    if st.button("PPT 생성"):
        prs = Presentation()
        for group in edited_grouped_slides:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            y_offset = 0.5
            for slide_content in group:
                title_shape = slide.shapes.add_textbox(Inches(0.7), Inches(y_offset), Inches(6.5), Inches(1))
                title_frame = title_shape.text_frame
                title_frame.clear()
                p_title = title_frame.add_paragraph()
                p_title.text = slide_content.get('title', '')
                p_title.font.size = Pt(36)
                p_title.font.bold = True
                p_title.alignment = PP_ALIGN.LEFT
                y_offset += 1
                content_shape = slide.shapes.add_textbox(Inches(0.7), Inches(y_offset), Inches(6.5), Inches(3.5))
                content_frame = content_shape.text_frame
                content_frame.clear()
                for line in slide_content.get('content', '').split('\n'):
                    if line.strip():
                        for wrapped_line in textwrap.wrap(line.strip(), width=40):
                            p = content_frame.add_paragraph()
                            p.text = wrapped_line
                            p.font.size = Pt(20)
                            p.alignment = PP_ALIGN.LEFT
                y_offset += 2.5
        ppt_io = io.BytesIO()
        prs.save(ppt_io)
        st.download_button("PPT 다운로드", ppt_io.getvalue(), file_name="output.pptx")
